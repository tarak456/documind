import os
import re

class DocumentProcessor:
    """Extracts plain text + metadata from PDF, DOCX, TXT, CSV, XLSX, PPTX."""

    def extract_text(self, file_path: str, original_name: str) -> tuple[str, dict]:
        ext = original_name.rsplit(".", 1)[-1].lower()

        extractors = {
            "pdf":  self._from_pdf,
            "docx": self._from_docx,
            "txt":  self._from_txt,
            "csv":  self._from_csv,
            "xlsx": self._from_xlsx,
            "xls":  self._from_xlsx,
            "pptx": self._from_pptx,
        }

        if ext not in extractors:
            raise ValueError(f"Unsupported file type: .{ext}")

        text = extractors[ext](file_path)
        text = self._clean(text)
        metadata = self._metadata(text, ext.upper())
        return text, metadata

    # ── Extractors ────────────────────────────────────────────────────────────

    def _from_pdf(self, path: str) -> str:
        import PyPDF2
        text_parts = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        if not text_parts:
            raise ValueError(
                "No text could be extracted from this PDF. "
                "It may be a scanned/image-based PDF."
            )
        return "\n\n".join(text_parts)

    def _from_docx(self, path: str) -> str:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Also grab table cells
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n\n".join(parts)

    def _from_txt(self, path: str) -> str:
        for enc in ("utf-8", "latin-1", "cp1252"):
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode the text file.")

    def _from_csv(self, path: str) -> str:
        import pandas as pd
        df = pd.read_csv(path, encoding_errors="replace")
        lines = [f"CSV Data — {len(df)} rows × {len(df.columns)} columns",
                 f"Columns: {', '.join(df.columns.tolist())}",
                 "", df.to_string(index=False, max_rows=500)]
        return "\n".join(lines)

    def _from_xlsx(self, path: str) -> str:
        import pandas as pd
        xl = pd.ExcelFile(path)
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            parts.append(
                f"Sheet: {sheet} ({len(df)} rows × {len(df.columns)} cols)\n"
                + df.to_string(index=False, max_rows=200)
            )
        return "\n\n".join(parts)

    def _from_pptx(self, path: str) -> str:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _clean(self, text: str) -> str:
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        return text.strip()

    def _metadata(self, text: str, file_type: str) -> dict:
        words = text.split()
        paras = [p for p in text.split("\n\n") if p.strip()]
        return {
            "word_count": len(words),
            "char_count": len(text),
            "para_count": len(paras),
            "file_type": file_type,
        }
